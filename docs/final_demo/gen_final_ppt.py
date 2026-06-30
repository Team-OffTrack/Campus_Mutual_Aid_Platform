#!/usr/bin/env python3
"""Generate final defense PPTX v4 (16:9, ~24 slides) for 校园互助服务平台.
   v4: stripped all theory explanations, enlarged fonts, fixed overflows,
       split testing tables, one screenshot per feature page, 4 section dividers.
   Requires: python-pptx, Pillow"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ═══════════════════════════════════════════════════════════════
# Color Palette
# ═══════════════════════════════════════════════════════════════
BG_DEEP    = RGBColor(0x0B, 0x11, 0x21)
BG_CARD    = RGBColor(0x16, 0x20, 0x36)
ACC_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)  # violet
ACC_CYAN   = RGBColor(0x22, 0xD3, 0xEE)  # cyan
ACC_GREEN  = RGBColor(0x10, 0xB9, 0x81)  # emerald
ACC_AMBER  = RGBColor(0xF5, 0x9E, 0x0B)  # amber
ACC_RED    = RGBColor(0xEF, 0x44, 0x44)  # red
ACC_PINK   = RGBColor(0xEC, 0x48, 0x99)  # pink
WHITE      = RGBColor(0xF8, 0xFA, 0xFC)
GRAY       = RGBColor(0x94, 0xA3, 0xB8)
GRAY_DARK  = RGBColor(0x64, 0x74, 0x8B)
GRAY_DIM   = RGBColor(0x47, 0x53, 0x69)

W = Inches(13.333)
H = Inches(7.5)
BASEDIR = os.path.dirname(os.path.abspath(__file__))
SCREENSHOT_DIR = os.path.join(BASEDIR, "screenshots")
FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ═══════════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════════

def set_bg(slide, color=BG_DEEP):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def tb(slide, left, top, width, height, text, *,
       size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = FONT
    p.alignment = align
    return box

def ml(slide, left, top, width, height, lines, *,
       size=14, color=GRAY, spacing=Pt(8)):
    """Multi-line text box."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, clr = item, False, color
        else:
            txt = item[0]; bld = item[1] if len(item)>1 else False
            clr = item[2] if len(item)>2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size); p.font.bold = bld
        p.font.color.rgb = clr; p.font.name = FONT
        p.space_after = spacing
    return box

def card(slide, left, top, width, height, color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def bar(slide, left, top, width, color=ACC_PURPLE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.05))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def stitle(slide, text, subtitle=None, accent=ACC_PURPLE):
    tb(slide, 0.9, 0.4, 11.5, 0.6, text, size=30, bold=True, color=WHITE)
    bar(slide, 0.9, 1.0, 1.8, accent)
    if subtitle:
        tb(slide, 0.9, 1.12, 11.5, 0.32, subtitle, size=13, color=GRAY)

def pn(slide, num):
    tb(slide, 12.3, 7.1, 0.7, 0.3, str(num), size=10, color=GRAY_DARK,
       align=PP_ALIGN.RIGHT)

def section_divider(slide, num, title, subtitle, accent=ACC_PURPLE):
    set_bg(slide)
    tb(slide, 0.9, 1.5, 3.0, 2.0, f"0{num}", size=96, bold=True,
       color=accent, align=PP_ALIGN.LEFT)
    bar(slide, 0.9, 3.6, 3.0, accent)
    tb(slide, 0.9, 3.9, 11.5, 1.0, title, size=42, bold=True, color=WHITE)
    tb(slide, 0.9, 4.9, 11.5, 0.5, subtitle, size=18, color=GRAY)
    # no page number on dividers

def screenshot(slide, left, top, width, height, name):
    """Insert screenshot or gray dashed placeholder."""
    path = os.path.join(SCREENSHOT_DIR, name)
    if os.path.exists(path):
        try:
            img = Image.open(path)
            iw, ih = img.size
            ratio = iw / ih; box_ratio = width / height
            if ratio > box_ratio:
                new_h = width / ratio
                slide.shapes.add_picture(path,
                    Inches(left), Inches(top + (height - new_h)/2),
                    Inches(width), Inches(new_h))
            else:
                new_w = height * ratio
                slide.shapes.add_picture(path,
                    Inches(left + (width - new_w)/2), Inches(top),
                    Inches(new_w), Inches(height))
            return
        except Exception:
            pass
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x12, 0x1A, 0x2E)
    shape.line.color.rgb = GRAY_DIM; shape.line.width = Pt(1.5)
    shape.line.dash_style = 2
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[ 截图占位: {name} ]\n请放入 docs/final_demo/screenshots/"
    p.font.size = Pt(11); p.font.color.rgb = GRAY_DARK
    p.font.name = FONT; p.alignment = PP_ALIGN.CENTER

def plantuml_img(slide, left, top, width, height, filename):
    """Insert PlantUML PNG with aspect ratio preserved."""
    path = os.path.join(BASEDIR, filename)
    if not os.path.exists(path):
        return tb(slide, left, top, width, height,
                  f"[ 缺失: {filename} ]", size=12, color=ACC_RED)
    img = Image.open(path); iw, ih = img.size
    ratio = iw / ih; box_ratio = width / height
    if ratio > box_ratio:
        new_h = width / ratio
        slide.shapes.add_picture(path,
            Inches(left), Inches(top + (height - new_h)/2),
            Inches(width), Inches(new_h))
    else:
        new_w = height * ratio
        slide.shapes.add_picture(path,
            Inches(left + (width - new_w)/2), Inches(top),
            Inches(new_w), Inches(height))

def add_table(slide, left, top, col_widths, headers, rows, *,
              header_size=14, cell_size=13, header_color=ACC_PURPLE):
    """Add formatted table."""
    n_rows = len(rows) + 1; n_cols = len(headers)
    tbl_w = sum(col_widths)
    row_h = 0.45
    tbl_shape = slide.shapes.add_table(n_rows, n_cols,
        Inches(left), Inches(top), Inches(tbl_w), Inches(row_h * n_rows))
    tbl = tbl_shape.table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci); cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(header_size); p.font.bold = True
            p.font.color.rgb = WHITE; p.font.name = FONT
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = header_color
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci); cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(cell_size); p.font.color.rgb = GRAY
                p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD if ri%2==0 else BG_DEEP
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tbl_shape

def hci_tags(slide, tags, y=6.85):
    """HCI keyword row at bottom of feature demo pages."""
    tb(slide, 0.6, y, 12.1, 0.4, "  ·  ".join(tags),
       size=10, color=GRAY_DIM, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — 封面
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
card(s, 0, 0, 13.333, 0.06, ACC_PURPLE)
tb(s, 1.2, 1.5, 10.9, 1.1, "校园互助服务平台", size=48, bold=True,
   color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 1.2, 2.6, 10.9, 0.6, "期末答辩", size=30, bold=True,
   color=ACC_CYAN, align=PP_ALIGN.CENTER)
tb(s, 1.2, 3.3, 10.9, 0.45, "软件工程与计算 II · 人机交互",
   size=17, color=GRAY, align=PP_ALIGN.CENTER)
card(s, 5.8, 3.9, 1.733, 0.013, ACC_PURPLE)
tb(s, 1.2, 4.2, 10.9, 0.45,
   "团队：不误正业     成员：侯乔岳 · 杨佳兴 · 沈诺 · 胡皓轩",
   size=15, color=GRAY, align=PP_ALIGN.CENTER)
tb(s, 1.2, 4.7, 10.9, 0.4, "沈诺      2026 年 6 月",
   size=13, color=GRAY_DARK, align=PP_ALIGN.CENTER)
pn(s, 1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — 团队介绍
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "团队介绍", "角色分工与职责")

add_table(s, 1.5, 1.6, [2.6, 2.0, 7.0],
    ["角  色", "成  员", "职  责"],
    [["需求负责人", "侯乔岳", "功能规划、需求文档"],
     ["架构负责人", "杨佳兴", "系统设计、数据库 Schema、API 设计"],
     ["开发负责人和测试负责人", "沈诺、胡皓轩",
      "前后端开发、代码实现、测试用例、质量保证"]])

members = [
    ("侯乔岳", "需求负责人", ACC_GREEN,
     ["功能规划与需求文档", "用户调研与用例编写",
      "验收标准制定", "需求优先级排序"]),
    ("杨佳兴", "架构负责人", ACC_AMBER,
     ["系统架构设计", "数据库 Schema 设计",
      "ADR 架构决策记录", "API 规范设计"]),
    ("沈诺", "开发 & 测试负责人", ACC_CYAN,
     ["前后端全栈开发", "243 测试用例编写",
      "CI/CD 流水线搭建", "HCI 交互细节落地"]),
    ("胡皓轩", "开发 & HCI 设计", ACC_PURPLE,
     ["需求功能设计", "HCI 交互设计",
      "M3E 设计系统", "演示与展示"]),
]
for i, (name, role, clr, items) in enumerate(members):
    x = 0.55 + i * 3.2
    card(s, x, 3.4, 2.95, 3.4)
    card(s, x, 3.4, 2.95, 0.06, clr)
    tb(s, x + 0.2, 3.6, 2.55, 0.55, name, size=20, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, x + 0.2, 4.15, 2.55, 0.32, role, size=12, color=clr,
       align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        tb(s, x + 0.2, 4.75 + j * 0.42, 2.55, 0.35,
           f"▸ {item}", size=12, color=GRAY)
pn(s, 2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — 项目概述
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "项目概述", "校园互助 — 跑腿 · 交易 · 组队 · 招领 · 互助 · 其他")

stats = [("15", "数据表"), ("57", "API 端点"), ("18", "前端路由"),
         ("243", "测试用例"), ("9", "成就徽章")]
for i, (num, label) in enumerate(stats):
    x = 0.8 + i * 2.45
    card(s, x, 1.5, 2.15, 1.25)
    tb(s, x, 1.55, 2.15, 0.6, num, size=30, bold=True,
       color=ACC_CYAN, align=PP_ALIGN.CENTER)
    tb(s, x, 2.15, 2.15, 0.3, label, size=13, color=GRAY,
       align=PP_ALIGN.CENTER)

types = [("跑腿代取", "代取快递/外卖", ACC_AMBER),
         ("二手交易", "旧书/电子产品", ACC_GREEN),
         ("组队匹配", "竞赛/学习搭子", ACC_CYAN),
         ("失物招领", "丢失/拾取物品", ACC_PURPLE),
         ("学习互助", "辅导/答疑", ACC_PINK),
         ("其   他", "搬家/装机等", GRAY)]
for i, (name, desc, clr) in enumerate(types):
    x = 0.8 + i * 2.1
    card(s, x, 3.1, 1.85, 1.0)
    card(s, x, 3.1, 1.85, 0.05, clr)
    tb(s, x + 0.1, 3.2, 1.65, 0.32, name, size=14, bold=True,
       color=clr, align=PP_ALIGN.CENTER)
    tb(s, x + 0.1, 3.55, 1.65, 0.28, desc, size=10, color=GRAY,
       align=PP_ALIGN.CENTER)

tb(s, 0.8, 4.45, 11.7, 0.3, "技术栈", size=17, bold=True, color=WHITE)
tb(s, 0.8, 4.85, 11.7, 0.28,
   "后端：Spring Boot 3.2.5 · MyBatis-Plus 3.5.6 · Spring Security + JWT · WebSocket STOMP · Flyway · MySQL",
   size=12, color=GRAY)
tb(s, 0.8, 5.2, 11.7, 0.28,
   "前端：Vue 3 (Composition API) · Vite 5 · Vant 4 · Pinia · Vue Router 4 · Axios · M3E 设计系统",
   size=12, color=GRAY)

features = ["JWT 无状态认证", "积分不可变账本", "WebSocket 实时通信",
            "9 种成就徽章", "响应式双视图", "举报审核系统",
            "悲观锁防超发", "Flyway 版本迁移", "统一 ApiResult"]
for i, f in enumerate(features):
    col = i % 3; row = i // 3
    tb(s, 0.8 + col * 4.1, 5.7 + row * 0.42, 3.9, 0.35,
       f"✦  {f}", size=13, color=GRAY)
pn(s, 3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — 敏捷开发
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "过程模型：敏捷开发", "Agile — 拥抱变化 · 短迭代 · 持续交付", ACC_CYAN)

# 4 Agile values
vals = ["个体和交互 > 过程和工具", "可工作的软件 > 面面俱到的文档",
        "客户合作 > 合同谈判", "响应变化 > 遵循计划"]
for i, v in enumerate(vals):
    x = 0.6 + i * 3.15
    card(s, x, 1.5, 2.95, 0.65)
    tb(s, x + 0.1, 1.55, 2.75, 0.55, v, size=13, bold=True,
       color=ACC_CYAN, align=PP_ALIGN.CENTER)

# P0-P4 phases
phases = [
    ("P0", "团队组建", "团队章程 · AI 协作规范 · 技术选型"),
    ("P1", "需求分析", "用户调研 · 用例建模 · 需求规格"),
    ("P2", "体系结构", "架构设计 ADR · 方案对比 · 辩论"),
    ("P3", "详细设计", "ER 图 · 类图 · API · DDL · SOLID"),
    ("P4", "编码测试", "Sprint 看板 · 编码实现 · CI/CD"),
]
for i, (phase, title, items) in enumerate(phases):
    x = 0.45 + i * 2.5
    card(s, x, 2.55, 2.3, 2.55)
    tb(s, x + 0.1, 2.65, 2.1, 0.5, phase, size=26, bold=True,
       color=ACC_PURPLE, align=PP_ALIGN.CENTER)
    tb(s, x + 0.1, 3.15, 2.1, 0.3, title, size=15, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, x + 0.1, 3.6, 2.1, 0.4, items, size=12, color=GRAY,
       align=PP_ALIGN.CENTER)
    if i < 4:
        tb(s, x + 2.3, 3.55, 0.2, 0.35, "›", size=22, bold=True,
           color=GRAY_DARK, align=PP_ALIGN.CENTER)

# Sprint strategy — simplified to one row of keywords
card(s, 0.8, 5.45, 11.7, 1.7)
tb(s, 1.0, 5.6, 11.3, 0.35, "Sprint 执行", size=17, bold=True, color=WHITE)
sprint_items = [
    "优先级驱动 P0→P4 — 核心功能先交付",
    "每轮迭代产出可运行增量（Demo-first）",
    "看板追踪：待办 / 进行中 / 已完成",
    "估时到小时 · 每日同步 · 回顾会议",
    "需求随迭代反馈灵活调整 — 拥抱变化",
]
for j, item in enumerate(sprint_items):
    col = j % 2; row = j // 2
    tb(s, 1.0 + col * 5.9, 6.05 + row * 0.38, 5.6, 0.32,
       f"• {item}", size=12, color=GRAY)
pn(s, 4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — 需求分析
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "需求分析", "需求三层次 · 六种需求类型 · 功能性 / 非功能性需求")

# Three levels
levels = [
    ("业务需求", "组织的高层目标",
     "解决校园互助信息不对称\n降低学生时间成本\n提升校园生活效率"),
    ("用户需求", "用户能用系统做什么",
     "发布需求、接单完成任务\n积分激励、实时沟通\n信用评价、成就激励"),
    ("系统需求", "系统必须提供的功能",
     "注册登录、需求 CRUD\n状态流转、积分账本\n通知推送、管理后台"),
]
for i, (level, desc, items) in enumerate(levels):
    x = 0.8 + i * 4.1
    card(s, x, 1.5, 3.85, 2.4)
    tb(s, x + 0.15, 1.6, 3.55, 0.35, level, size=19, bold=True,
       color=ACC_CYAN, align=PP_ALIGN.CENTER)
    tb(s, x + 0.15, 2.0, 3.55, 0.28, desc, size=12, color=GRAY,
       align=PP_ALIGN.CENTER)
    ml(s, x + 0.2, 2.45, 3.35, 1.3, items.split("\n"), size=13, spacing=Pt(10))

# 6 types + functional/non-functional
tb(s, 0.8, 4.2, 11.7, 0.32, "六种需求类型", size=17, bold=True, color=WHITE)
type_names = ["跑腿代取", "二手交易", "组队匹配", "失物招领", "学习互助", "其他"]
for i, name in enumerate(type_names):
    card(s, 0.8 + i * 2.1, 4.65, 1.85, 0.55)
    tb(s, 0.8 + i * 2.1, 4.68, 1.85, 0.48, name, size=14, bold=True,
       color=ACC_CYAN, align=PP_ALIGN.CENTER)

# Functional / Non-functional (compact, one line each)
tb(s, 0.8, 5.55, 11.7, 0.28, "功能性需求", size=14, bold=True, color=ACC_CYAN)
tb(s, 0.8, 5.88, 11.7, 0.28,
   "注册登录 · 需求发布/接单/完成 · 实时聊天 · 积分系统 · 评价评分 · 成就徽章 · 举报审核 · 管理仪表盘",
   size=12, color=GRAY)
tb(s, 0.8, 6.3, 11.7, 0.28, "非功能性需求", size=14, bold=True, color=ACC_PURPLE)
tb(s, 0.8, 6.63, 11.7, 0.28,
   "JWT 鉴权 · 角色守卫 · 悲观锁防超发 · 积分不可变账本 · 响应式设计 · 键盘可访问 · 防重复提交",
   size=12, color=GRAY)
pn(s, 5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Section Divider: 系统架构
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, 2, "系统架构",
    "体系结构设计 · 包结构 · 数据库设计", ACC_CYAN)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — 体系结构设计
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "体系结构设计", "分层风格 · 模块化单体 · 关键架构决策", ACC_CYAN)

# Architecture diagram — left 58%
# Left: Three-layer text listing (replaces architecture diagram)
card(s, 0.3, 1.4, 4.0, 5.6)
tb(s, 0.5, 1.5, 3.6, 0.35, "分层架构（Layer Style）", size=16, bold=True, color=ACC_CYAN)

layers_data = [
    ("展示层", ACC_PURPLE,
     "Vue 3 · Vite 5 · Vant 4\nPinia · Vue Router 4 · Axios\nM3E 设计系统\n响应式双视图 @media 768px"),
    ("逻辑层", ACC_CYAN,
     "Spring Boot 3.2.5\nSpring Security + JWT\nMyBatis-Plus 3.5.6\nWebSocket STOMP"),
    ("数据层", ACC_GREEN,
     "MySQL · Flyway V1→V14\n单表继承 discriminator\n积分不可变账本\n悲观锁 SELECT FOR UPDATE"),
]
for i, (name, clr, items) in enumerate(layers_data):
    y = 2.0 + i * 1.7
    card(s, 0.55, y, 3.55, 1.55)
    card(s, 0.55, y, 3.55, 0.05, clr)
    tb(s, 0.75, y + 0.1, 3.15, 0.35, name, size=15, bold=True, color=clr)
    ml(s, 0.75, y + 0.5, 3.15, 0.95, items.split("\n"), size=12, spacing=Pt(6))

# Vertical connectors between layers
for i in range(2):
    y_mid = 3.55 + i * 1.7
    tb(s, 2.1, y_mid + 0.05, 0.3, 0.25, "▼", size=12, color=GRAY_DIM, align=PP_ALIGN.CENTER)

# ADR cards — right 48%
tb(s, 8.3, 1.4, 4.6, 0.35, "关键架构决策 (ADR)", size=17, bold=True, color=ACC_CYAN)
adrs = [
    ("模块化单体 > 微服务",
     "10 周交付 · 无分布式事务 · 模块边界清晰可后续拆分"),
    ("JWT 无状态 > Session",
     "跨域灵活 · 无需服务端存储 · 适合前后端分离"),
    ("单表继承 > 多表继承",
     "discriminator + JSON 列 · 6 种需求类型共享 demand 表"),
    ("悲观锁 > 乐观锁",
     "积分强一致性 · SELECT FOR UPDATE · 防止并发超发"),
]
for i, (title, reason) in enumerate(adrs):
    y = 2.0 + i * 1.2
    card(s, 8.3, y, 4.5, 1.05)
    tb(s, 8.5, y + 0.08, 4.1, 0.35, title, size=11, bold=True, color=ACC_CYAN)
    tb(s, 8.5, y + 0.43, 4.1, 0.52, reason, size=10, color=GRAY)
pn(s, 7)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — 包结构（单独一页，图铺满）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "包结构", "逻辑包图", ACC_CYAN)

plantuml_img(s, 0.3, 1.2, 12.7, 5.7, "packages.png")

tb(s, 0.8, 7.05, 11.7, 0.28,
   "controller 10 · service 14 · mapper 17 · entity 16 · dto 12 · config 7 · security 3 · common 5",
   size=12, color=GRAY, align=PP_ALIGN.CENTER)
pn(s, 8)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — 数据库设计（大字关键词，单独一页）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "数据库设计", "15 张表 · Flyway V1→V14 · 关键设计决策", ACC_CYAN)

db_keywords = [
    ("15 张表", "Flyway V1 → V14\n版本化增量迁移"),
    ("单表继承", "discriminator 字段\n6 种需求类型共用 demand 表"),
    ("不可变账本", "points_transaction 仅 INSERT\nbalance_after 快照审计"),
    ("悲观锁", "SELECT … FOR UPDATE\n防止积分并发超发"),
    ("UNIQUE 防重", "签到日期 · 收藏对\n会话 ID · 学号"),
]
for i, (title, desc) in enumerate(db_keywords):
    x = 0.5 + i * 2.5
    card(s, x, 1.5, 2.3, 1.65)
    card(s, x, 1.5, 2.3, 0.05, ACC_CYAN)
    tb(s, x + 0.12, 1.6, 2.06, 0.42, title, size=17, bold=True,
       color=ACC_CYAN, align=PP_ALIGN.CENTER)
    ml(s, x + 0.12, 2.15, 2.06, 0.9, desc.split("\n"), size=12)

# 15 tables list (big text)
tb(s, 0.8, 3.6, 11.7, 0.35, "15 张数据表", size=18, bold=True, color=WHITE)
tables_list = [
    "user  ·  user_account  ·  demand  ·  demand_application  ·  demand_favorite",
    "points_transaction  ·  badge_definition  ·  user_badge  ·  evaluation  ·  report",
    "chat_session  ·  chat_message  ·  notification  ·  team  ·  team_member",
]
for i, line in enumerate(tables_list):
    tb(s, 0.8, 4.15 + i * 0.55, 11.7, 0.45, line,
       size=16, color=GRAY, align=PP_ALIGN.CENTER)
pn(s, 9)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Section Divider: 开发实现
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, 3, "开发实现：功能与 HCI",
    "API 规范 · 动态建模 · 功能演示 · 人机交互", ACC_PURPLE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — API 规范 + 类图
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "API 规范与设计类图", "57 端点 · 10 Controller · 统一响应")

# Left: API keywords
card(s, 0.5, 1.35, 4.3, 2.6)
tb(s, 0.7, 1.45, 3.9, 0.35, "REST API 规范", size=17, bold=True, color=ACC_PURPLE)
ml(s, 0.7, 1.9, 3.9, 1.9, [
    ("57 REST 端点 + 4 WebSocket", True, WHITE),
    "",
    "10 个 Controller：",
    "User · Demand · Chat · Badge",
    "Evaluation · Notification · Points",
    "Report · TeamMember · Admin",
    "",
    ("统一响应：{ code, message, data }", False, ACC_CYAN),
    ("JWT 鉴权 + 角色守卫 USER/ADMIN", False, ACC_CYAN),
    ("全局异常处理 @RestControllerAdvice", False, ACC_CYAN),
], size=12)

# Right: class diagram (big — 62% of page)
plantuml_img(s, 5.0, 1.35, 8.0, 5.9, "class_diagram.png")
pn(s, 11)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — 状态图（图铺满）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "Demand 需求状态图", "State Diagram")

plantuml_img(s, 0.3, 1.25, 12.7, 5.6, "state_diagram.png")

tb(s, 0.8, 7.0, 11.7, 0.3,
   "OPEN → IN_PROGRESS → COMPLETED / CANCELLED",
   size=14, color=GRAY, align=PP_ALIGN.CENTER)
pn(s, 12)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — 顺序图（图铺满）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "系统顺序图：「发布跑腿需求」", "SSD — 展示层 → 逻辑层 → 数据层")

plantuml_img(s, 0.3, 1.25, 12.7, 5.6, "sequence.png")

tb(s, 0.8, 7.0, 11.7, 0.3,
   "POST /api/v1/demands → 3 条 SQL 原子事务 → ApiResult → Toast",
   size=14, color=GRAY, align=PP_ALIGN.CENTER)
pn(s, 13)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — 登录与首页（全页截图）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "功能演示：登录与首页", "")

screenshot(s, 0.6, 1.2, 12.1, 5.4, "login.png")

hci_tags(s, [
    "EEC 闭环（7 阶段）",
    "Nielsen #1 系统状态的可见度 — 签到两态切换",
    "Shneiderman #3 信息丰富的反馈 — 焦点光晕 + toast",
    "问候语按时间切换 — UX 有温度",
])
pn(s, 14)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — 需求广场（全页截图）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "功能演示：需求广场", "")

screenshot(s, 0.6, 1.2, 12.1, 5.4, "demand_list_desktop.png")

hci_tags(s, [
    "响应式 @media 768px — 卡片 ↔ 表格",
    "Fitts 定律 — FAB 拇指热区 + spring 回弹",
    "Nielsen #6 识别优于记忆 — 彩色类型 chip + 图标",
    "Shneiderman #2 符合普遍可用性",
])
pn(s, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — 发布需求（全页截图）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "功能演示：发布需求", "")

screenshot(s, 0.6, 1.2, 12.1, 5.4, "demand_publish.png")

hci_tags(s, [
    "Shneiderman #5 预防并处理错误 — 按钮变灰防重复提交",
    "Shneiderman #4 结束信息 — 三重反馈闭环（Toast + 跳转 + 列表更新）",
    "Spring 弹簧动画 — 类型 chip 点击物理手感",
    "必填项缺失 — 红色实时校验提示",
])
pn(s, 16)

# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — 需求详情（全页截图）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "功能演示：需求详情", "")

screenshot(s, 0.6, 1.2, 12.1, 5.4, "demand_detail.png")

hci_tags(s, [
    "Shneiderman #7 支持内部控制点 — 不可逆操作二次确认",
    "Shneiderman #5 预防错误 — 接单/完成/取消均需确认弹框",
    "Shneiderman #6 让操作容易撤销 — 评分提交后可修改",
    "Nielsen #1 状态颜色语义 — OPEN 蓝 · 进行中 琥珀 · 完成 绿 · 取消 红灰",
])
pn(s, 17)

# ═══════════════════════════════════════════════════════════════
# SLIDE 18 — 聊天 & 徽章（左右并排两张截图）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "功能演示：私信聊天 & 成就徽章", "")

screenshot(s, 0.5, 1.25, 6.1, 3.5, "chat.png")
screenshot(s, 6.8, 1.25, 6.1, 3.5, "badges.png")

hci_tags(s, [
    "格式塔 — 图底原则（不对称气泡圆角）",
    "WebSocket 实时推送 — Nielsen #1 最高标准",
    "Peak-End Rule — 徽章全屏动效 spring(0→1.1→1.0) + teleport",
    "游戏化激励 — 佩戴角标全局可见 · 社交展示",
])
pn(s, 18)

# ═══════════════════════════════════════════════════════════════
# SLIDE 19 — 管理后台（全页截图）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "功能演示：管理后台", "")

screenshot(s, 0.6, 1.2, 12.1, 5.4, "admin_dashboard.png")

hci_tags(s, [
    "Nielsen #8 审美感和最小化设计 — 4 统计卡片信息架构",
    "Nielsen #1 状态可见度 — 举报红色角标 + 高亮双重强调",
    "数据格式化 1.2k / 1.2w — 减轻认知负担",
    "需求类型分布 chips + 百分比 — 数据图形化 > 纯数字",
])
pn(s, 19)

# ═══════════════════════════════════════════════════════════════
# SLIDE 20 — Section Divider: 测试
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, 4, "测试",
    "白盒测试 · 黑盒测试 · CI/CD", ACC_GREEN)

# ═══════════════════════════════════════════════════════════════
# SLIDE 21 — 白盒测试（独立页，两张表，14pt+）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "白盒测试", "三种覆盖标准：语句覆盖 · 分支覆盖 · 路径覆盖", ACC_GREEN)

# Table 1: three coverage standards
card(s, 0.7, 1.5, 11.9, 2.5)
tb(s, 0.9, 1.6, 11.5, 0.35, "覆盖标准（从弱到强）", size=17, bold=True, color=WHITE)
add_table(s, 0.9, 2.05, [2.0, 4.5, 1.8, 3.6],
    ["标  准", "要  求", "强  度", "最少用例数"],
    [["语句覆盖", "每条语句至少执行一次", "★ 最弱", "可能 1 个用例"],
     ["分支覆盖", "每个判断的真 / 假分支各执行一次", "★★ 中等", "每个 if 需要 ≥ 2 个 (T+F)"],
     ["路径覆盖", "每条独立执行路径各一次", "★★★ 最强",
      "路径数指数级\n不可行时用 V(G) 代替"]],
    header_size=15, cell_size=14, header_color=ACC_GREEN)

# Table 2: project practice
card(s, 0.7, 4.25, 11.9, 3.05)
tb(s, 0.9, 4.35, 11.5, 0.35, "本项目测试实践", size=17, bold=True, color=WHITE)
add_table(s, 0.9, 4.8, [2.8, 3.5, 1.8, 3.8],
    ["测试层次", "框架 & 工具", "用例数", "测试重点"],
    [["Service 层（白盒）", "JUnit 5 + @SpringBootTest\nH2 内存数据库", "135",
      "语句：register → 验证字段\n分支：重复学号 → assertThrows\n路径：发布→接单→完成→互评"],
     ["Controller 层\n（白盒 + 黑盒）", "MockMvc\n@AutoConfigureMockMvc", "102",
      "等价类：200 / 400 / 403\n边界值：min=1 合法\n多角色：发布者/接单者/管理员"],
     ["WebSocket 集成", "STOMP Client\nBlockingQueue", "6",
      "JWT WebSocket 鉴权\n消息实时推送验证"]],
    header_size=15, cell_size=14, header_color=ACC_GREEN)
pn(s, 21)

# ═══════════════════════════════════════════════════════════════
# SLIDE 22 — 黑盒测试：等价类 + 边界值（独立页，两个表并排，14pt+）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "黑盒测试", "等价类划分 · 边界值分析", ACC_GREEN)

# Left: Equivalence class
tb(s, 0.7, 1.5, 5.8, 0.35, "等价类划分", size=18, bold=True, color=WHITE)
card(s, 0.7, 1.95, 6.0, 2.65)
add_table(s, 0.9, 2.1, [1.8, 2.2, 1.8],
    ["等价类类型", "如何确定", "示例（积分 1–100）"],
    [["有效等价类", "符合需求范围", "1 ≤ x ≤ 100\n取代表值 50"],
     ["无效等价类", "超出范围", "x < 1（取 0）\nx > 100（取 101）"],
     ["额外无效类", "非预期类型", "字母 · 空值\n负数 · 小数"]],
    header_size=15, cell_size=14, header_color=ACC_GREEN)

# Right: Boundary value
tb(s, 7.2, 1.5, 5.8, 0.35, "边界值分析 — 最易出 Bug！", size=18, bold=True, color=WHITE)
card(s, 7.2, 1.95, 5.8, 2.65)
add_table(s, 7.4, 2.1, [1.4, 1.6, 1.4, 1.2],
    ["取值点", "公式", "积分(1–100)", "判定"],
    [["min−1", "min − 1", "0", "非法"],
     ["min", "min", "1", "合法"],
     ["min+1", "min + 1", "2", "合法"],
     ["max−1", "max − 1", "99", "合法"],
     ["max", "max", "100", "合法"],
     ["max+1", "max + 1", "101", "非法"]],
    header_size=15, cell_size=14, header_color=ACC_GREEN)

# One-line summary
tb(s, 0.7, 4.85, 11.9, 1.5,
   "从需求规格中提取：等价类划分 → 每个类选代表值\n边界值分析 → 对每个范围取 min−1 / min / min+1 / max−1 / max / max+1 六点",
   size=16, color=GRAY, align=PP_ALIGN.CENTER)
pn(s, 22)

# ═══════════════════════════════════════════════════════════════
# SLIDE 23 — CI/CD（独立页）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "CI/CD 持续集成与部署", "GitHub Actions · 自动化测试 · 持续交付")

# Backend CI — bigger, left half
card(s, 0.5, 1.4, 6.15, 3.8)
tb(s, 0.7, 1.5, 5.75, 0.4, "Backend CI", size=18, bold=True, color=ACC_CYAN)
ml(s, 0.7, 2.05, 5.75, 2.8, [
    ("触发：每次 push / Pull Request", False, GRAY),
    "",
    "① Checkout 代码 + Set up JDK 17",
    "② Maven 依赖缓存（actions/cache）",
    "③ mvn test（H2 内存数据库）",
    "④ 243 测试用例自动执行",
    "⑤ Checkstyle 代码风格检查",
    "",
    ("⏱ 平均耗时 ~1m20s", True, ACC_CYAN),
    ("✓ 全绿才能合并到 main", True, ACC_CYAN),
], size=14, spacing=Pt(6))

# Frontend CI — right half
card(s, 6.95, 1.4, 6.15, 3.8)
tb(s, 7.15, 1.5, 5.75, 0.4, "Frontend CI", size=18, bold=True, color=ACC_PURPLE)
ml(s, 7.15, 2.05, 5.75, 2.8, [
    ("触发：每次 push / Pull Request", False, GRAY),
    "",
    "① Checkout 代码 + Set up Node.js",
    "② npm ci（锁定依赖版本）",
    "③ npm run build（Vite 生产构建）",
    "④ ESLint 代码规范检查",
    "⑤ 构建产物验证（dist/ 非空）",
    "",
    ("⏱ 平均耗时 ~24s", True, ACC_PURPLE),
    ("✓ 构建失败阻止合并", True, ACC_PURPLE),
], size=14, spacing=Pt(6))

# Deploy info
card(s, 0.5, 5.55, 12.3, 1.55)
tb(s, 0.7, 5.65, 11.9, 0.35, "部署架构与分支策略", size=16, bold=True, color=WHITE)
ml(s, 0.7, 6.1, 11.9, 0.85, [
    ("部署架构", True, ACC_CYAN),
    "Nginx 反向代理 → HTTPS 统一入口 → /api/* 转发 Spring Boot :8080 · 其余转发 Vite :5173",
    "JWT 无状态认证 → 水平扩展友好 · 自签名证书（开发）/ Let's Encrypt（生产）",
    ("分支策略", True, ACC_CYAN),
    "main 保护分支 + feature/* 特性分支 · PR + Code Review → Squash Merge · 合并前强制通过 CI 检查",
], size=13, spacing=Pt(4))
pn(s, 23)

# ═══════════════════════════════════════════════════════════════
# SLIDE 24 — 总结
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
stitle(s, "项目总结", "功能完整度 · 工程规范 · 人机交互")

cols = [
    ("功能完整度", ACC_CYAN, [
        "6 种需求类型全覆盖",
        "实时聊天 WebSocket",
        "积分系统（不可变账本）",
        "9 种成就徽章 + 游戏化",
        "组队协作 · 收藏 · 举报审核",
    ]),
    ("工程规范", ACC_PURPLE, [
        "全链路文档 P0–P4",
        "Flyway V1–V14 版本迁移",
        "架构决策记录 ADR",
        "57 API + 18 路由",
        "243 测试 · 60%+ 覆盖率",
    ]),
    ("人机交互", ACC_GREEN, [
        "Nielsen 10 条全覆盖",
        "Shneiderman 8 条全覆盖",
        "EEC 模型 · Fitts 定律",
        "格式塔 · Peak-End Rule",
        "M3E 设计系统 200+ token",
    ]),
]
for i, (title, clr, items) in enumerate(cols):
    x = 0.7 + i * 4.15
    card(s, x, 1.5, 3.85, 3.3)
    card(s, x, 1.5, 3.85, 0.05, clr)
    tb(s, x + 0.15, 1.65, 3.55, 0.45, title, size=20, bold=True,
       color=clr, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        tb(s, x + 0.2, 2.4 + j * 0.55, 3.45, 0.45,
           f"✦  {item}", size=14, color=GRAY)

card(s, 0.5, 5.15, 12.3, 0.7)
tb(s, 0.7, 5.3, 11.9, 0.4,
   "GitHub Actions 双 CI 流水线  ·  243 测试全部通过  ·  文档完整可交付  ·  Nielsen + Shneiderman 全覆盖",
   size=14, bold=True, color=ACC_CYAN, align=PP_ALIGN.CENTER)
pn(s, 24)

# ═══════════════════════════════════════════════════════════════
# SLIDE 25 — 谢谢老师（独立页）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
tb(s, 1.2, 2.2, 10.9, 1.2, "谢谢老师！", size=52, bold=True,
   color=WHITE, align=PP_ALIGN.CENTER)
card(s, 5.0, 3.5, 3.333, 0.013, ACC_PURPLE)
tb(s, 1.2, 4.0, 10.9, 0.6, "欢迎提问 🙋", size=24, color=GRAY,
   align=PP_ALIGN.CENTER)
tb(s, 1.2, 5.0, 10.9, 0.5,
   "校园互助服务平台  ·  团队不误正业  ·  2026.06",
   size=14, color=GRAY_DARK, align=PP_ALIGN.CENTER)
pn(s, 25)

# ═══════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════
out = os.path.join(BASEDIR, "期末答辩.pptx")
prs.save(out)
print(f"✓ Saved: {out}")
print(f"  Size: {os.path.getsize(out):,} bytes")
print(f"  Slides: {len(prs.slides)}")
