#!/usr/bin/env python3
"""Generate the HCI member contribution PPTX (one slide, 16:9).
   Requires: python-pptx (pacman -S python-pptx)"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# --- title ---
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.4), Inches(11.333), Inches(0.8))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "校园互助服务平台 — 成员分工与贡献比例"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

sub_box = slide.shapes.add_textbox(Inches(1), Inches(1.1), Inches(11.333), Inches(0.5))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "人机交互课程 — 个人工作与贡献说明"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0xB0, 0xB0, 0xCC)
p.alignment = PP_ALIGN.CENTER

members = [
    {
        "name": "沈诺（组长）",
        "role": "前后端开发 · 测试与质量保证 · HCI 交互细节讨论",
        "detail": [
            "实现全部前端页面（Vue 3 + Vant 4，18 条路由）及可复用组件",
            "实现全部后端 API（Spring Boot，57 个端点）及 Service 业务逻辑",
            "编写 246 个后端测试用例，配置 CI lint 流水线",
            "与胡皓轩讨论确定 HCI 交互细节并落地为代码",
            "设计并实现成就徽章系统（9 种徽章自动检测 + 全屏动效）",
            "管理后台（仪表盘/用户/需求/举报管理）前后端实现",
        ],
        "pct": "50%",
        "color": RGBColor(0x06, 0xB6, 0xD4),
    },
    {
        "name": "胡皓轩",
        "role": "需求与功能设计 · 人机交互设计 · 项目展示与答辩",
        "detail": [
            "主导全部六种需求类型的功能规划与交互流程设计",
            "设计 M3E 设计系统（200+ CSS token）及响应式断点方案",
            "设计徽章弹窗动效、管理仪表盘信息架构等 HCI 细节",
            "与沈诺讨论确定 HCI 交互方案的代码实现",
            "参与前端页面开发与后端接口联调测试",
            "负责所有阶段性答辩与线下演示",
        ],
        "pct": "50%",
        "color": RGBColor(0x7C, 0x3A, 0xED),
    },
]

card_w = Inches(5.2)
card_h = Inches(5.0)
gap = Inches(0.8)
start_x = (prs.slide_width - (2 * card_w + gap)) / 2
start_y = Inches(1.8)

for i, m in enumerate(members):
    x = start_x + i * (card_w + gap)
    y = start_y

    shape = slide.shapes.add_shape(1, x, y, card_w, card_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
    shape.line.fill.background()

    bar = slide.shapes.add_shape(1, x, y, card_w, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = m["color"]
    bar.line.fill.background()

    name_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), card_w - Inches(0.6), Inches(0.55))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = m["name"]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    role_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.88), card_w - Inches(0.6), Inches(0.55))
    tf = role_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = m["role"]
    p.font.size = Pt(11)
    p.font.color.rgb = m["color"]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    div = slide.shapes.add_shape(1, x + Inches(0.8), y + Inches(1.55), card_w - Inches(1.6), Inches(0.015))
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor(0x40, 0x40, 0x60)
    div.line.fill.background()

    detail_box = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(1.75), card_w - Inches(0.5), Inches(2.5))
    tf = detail_box.text_frame
    tf.word_wrap = True
    for j, d in enumerate(m["detail"]):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + d
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xDD)
        p.space_after = Pt(5)
        p.line_spacing = Pt(14.5)

    pct_box = slide.shapes.add_textbox(x + Inches(0.5), y + card_h - Inches(0.65), card_w - Inches(1.0), Inches(0.45))
    tf = pct_box.text_frame
    p = tf.paragraphs[0]
    p.text = "贡献占比  " + m["pct"]
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0xAA)
    p.alignment = PP_ALIGN.CENTER

# --- bottom note ---
note_box = slide.shapes.add_textbox(Inches(1.5), Inches(7.0), Inches(10.333), Inches(0.4))
tf = note_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "贡献占比：沈诺与胡皓轩基本相同（各约 50%），共同负责交互设计决策的讨论与落地。项目另两位成员（侯乔岳、杨佳兴）未选修本课程，在软工项目中提供了需求分析与架构设计的少量帮助。"
p.font.size = Pt(10)
p.font.color.rgb = RGBColor(0x77, 0x77, 0x88)
p.alignment = PP_ALIGN.CENTER

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "成员分工与贡献比例.pptx")
prs.save(out)
print(f"Saved: {out} ({os.path.getsize(out)} bytes)")
